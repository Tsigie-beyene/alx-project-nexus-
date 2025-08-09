from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme colors
COLOR_BG = RGBColor(11, 16, 32)        # #0B1020 dark navy
COLOR_TEXT = RGBColor(245, 247, 250)   # near white
COLOR_MUTED = RGBColor(160, 174, 192)  # gray
COLOR_ACCENT = RGBColor(108, 92, 231)  # #6C5CE7 purple
COLOR_ACCENT_2 = RGBColor(0, 209, 178) # #00D1B2 teal

REPO_URL = "https://github.com/Tsigie-beyene/alx-project-nexus-"
PROJECT_NAME = "Project Nexus — E‑Commerce Backend (ProDev BE)"
AUTHOR = "[Your Name]"


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text, page_number):
    # Use the slide's presentation dimensions when available; fall back to defaults
    try:
        prs = slide.part.presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    except Exception:
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)

    textbox = slide.shapes.add_textbox(Inches(0.5), slide_height - Inches(0.6), slide_width - Inches(1.0), Inches(0.4))
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{text}"
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"{page_number}"
    r2.font.size = Pt(10)
    r2.font.color.rgb = COLOR_MUTED
    p2.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title, subtitle, links_line):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, COLOR_BG)

    # Accent bar
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.25)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLOR_ACCENT
    slide.shapes[-1].line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), prs.slide_width - Inches(1.6), Inches(1.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = COLOR_TEXT

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), prs.slide_width - Inches(1.6), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    r2.font.color.rgb = COLOR_MUTED

    # Links line
    link_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), prs.slide_width - Inches(1.6), Inches(0.8))
    tf3 = link_box.text_frame
    tf3.clear()
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = links_line
    r3.font.size = Pt(14)
    r3.font.color.rgb = COLOR_ACCENT_2

    add_footer(slide, PROJECT_NAME, 1)


def add_section_slide(prs, section_title, section_subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)

    # Accent block
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(0.15), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.5), prs.slide_width - Inches(2.0), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section_title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = COLOR_TEXT

    if section_subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), prs.slide_width - Inches(2.0), Inches(0.8))
        tf2 = sub_box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = section_subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = COLOR_MUTED

    return slide


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width - Inches(1.6), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLOR_TEXT

    # Bullet list
    y = Inches(1.8)
    content_box = slide.shapes.add_textbox(Inches(0.8), y, prs.slide_width - Inches(1.6), prs.slide_height - y - Inches(1.2))
    t = content_box.text_frame
    t.clear()

    for idx, line in enumerate(bullets):
        if idx == 0:
            p = t.paragraphs[0]
        else:
            p = t.add_paragraph()
            p.level = 0
        run = p.add_run()
        run.text = f"• {line}"
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_TEXT

    return slide


def add_two_column_slide(prs, title, left_bullets, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width - Inches(1.6), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLOR_TEXT

    # Columns
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5))
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))

    for box, bullets in [(left_box, left_bullets), (right_box, right_bullets)]:
        tf_box = box.text_frame
        tf_box.clear()
        for idx, line in enumerate(bullets):
            if idx == 0:
                p = tf_box.paragraphs[0]
            else:
                p = tf_box.add_paragraph()
            run = p.add_run()
            run.text = f"• {line}"
            run.font.size = Pt(18)
            run.font.color.rgb = COLOR_TEXT

    return slide


def add_api_examples_slide(prs):
    bullets = [
        "Base: /api/v1/",
        "Auth: POST /auth/login → { access, refresh }",
        "Categories: GET/POST /categories/, GET/PATCH/DELETE /categories/{id}/",
        "Products: GET/POST /products/, GET/PATCH/DELETE /products/{id}/",
        "List params: ?category=<id>&ordering=-price&page=1&page_size=12",
    ]
    return add_bullets_slide(prs, "API Design (REST)", bullets)


def add_erd_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_BG)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width - Inches(1.6), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "ERD — Core Data Model"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLOR_TEXT

    # Entity boxes
    box_w, box_h = Inches(3.5), Inches(2.0)
    user_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), box_w, box_h)
    cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.0), box_w, box_h)
    prod_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.0), box_w, box_h)

    for box in [user_box, cat_box, prod_box]:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(23, 29, 50)
        box.line.color.rgb = COLOR_ACCENT

    # Box titles and fields
    def fill_box_text(shape, title, fields):
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = COLOR_ACCENT_2
        for f in fields:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = f"- {f}"
            run2.font.size = Pt(14)
            run2.font.color.rgb = COLOR_TEXT

    fill_box_text(user_box, "User", [
        "id (PK)", "email (unique)", "password_hash", "is_staff",
    ])
    fill_box_text(cat_box, "Category", [
        "id (PK)", "name (unique)", "slug (unique)",
    ])
    fill_box_text(prod_box, "Product", [
        "id (PK)", "name", "description", "price", "stock",
        "category_id (FK)", "created_at", "updated_at",
    ])

    # Legend / relationships
    legend = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(7.0), Inches(1.2))
    tlf = legend.text_frame
    tlf.clear()
    p = tlf.paragraphs[0]
    r = p.add_run()
    r.text = "Relationships: Category 1—N Product | User used for auth/admin"
    r.font.size = Pt(16)
    r.font.color.rgb = COLOR_MUTED

    return slide


def build_presentation(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1) Title
    add_title_slide(
        prs,
        title=PROJECT_NAME,
        subtitle=f"Author: {AUTHOR}  |  GitHub: {REPO_URL}",
        links_line="Links: Repo • Live API • Swagger • ERD • Demo Video",
    )

    # 2) Problem & Goal
    s = add_bullets_slide(prs, "Problem & Goal", [
        "Problem: Power product discovery and management at scale",
        "Users: Shoppers, Admins",
        "Goal: Secure, performant, documented backend with clean APIs",
    ])
    add_footer(s, PROJECT_NAME, 2)

    # 3) Solution Overview
    s = add_bullets_slide(prs, "Solution Overview", [
        "CRUD: products, categories, users (JWT auth)",
        "Discovery: filtering, sorting, pagination",
        "Quality: indexing, query optimization, OpenAPI docs",
    ])
    add_footer(s, PROJECT_NAME, 3)

    # 4) Tech Stack & Tools
    s = add_two_column_slide(
        prs,
        "Tech Stack & Tools",
        [
            "Backend: Django, DRF",
            "DB: PostgreSQL",
            "Auth: JWT (access/refresh)",
            "Docs: Swagger/OpenAPI",
        ],
        [
            "DevEx: pre-commit, Black, Flake8, isort",
            "CI/CD: GitHub Actions",
            "Hosting: Railway/Render/Heroku/VPS",
            "Observability: Admin, Logs, Health",
        ],
    )
    add_footer(s, PROJECT_NAME, 4)

    # 5) Architecture
    s = add_bullets_slide(prs, "Architecture (High Level)", [
        "Client → Django/DRF → PostgreSQL",
        "Optional: Redis for caching/rate limiting",
        "Static/admin: Django Admin",
        "CI: test → lint → migrate → deploy",
    ])
    add_footer(s, PROJECT_NAME, 5)

    # 6) ERD
    s = add_erd_slide(prs)
    add_footer(s, PROJECT_NAME, 6)

    # 7) Database Design Decisions
    s = add_bullets_slide(prs, "Database Design Decisions", [
        "Normalization: 3NF; product belongs to one category",
        "Constraints: unique (category.name, category.slug), NOT NULL on critical fields",
        "Indexes: product(category_id, price, created_at)",
        "Optional: trigram/GIN for search on product.name",
    ])
    add_footer(s, PROJECT_NAME, 7)

    # 8) API Design
    s = add_api_examples_slide(prs)
    add_footer(s, PROJECT_NAME, 8)

    # 9) Filtering, Sorting, Pagination
    s = add_bullets_slide(prs, "Filtering, Sorting, Pagination", [
        "Filters: ?category=<id>",
        "Sorting: ?ordering=price or ?ordering=-price",
        "Pagination: ?page=2&page_size=20 (defaults and max limits)",
        "Example: /api/v1/products?category=3&ordering=-price&page=1&page_size=12",
    ])
    add_footer(s, PROJECT_NAME, 9)

    # 10) Authentication & Authorization
    s = add_bullets_slide(prs, "Authentication & Authorization", [
        "JWT: login returns access + refresh",
        "Password hashing: PBKDF2/Argon2",
        "Permissions: staff/admin required for write; read open",
        "Optional: throttling for sensitive endpoints",
    ])
    add_footer(s, PROJECT_NAME, 10)

    # 11) Performance & Optimization
    s = add_bullets_slide(prs, "Performance & Optimization", [
        "ORM: select_related('category') for product lists",
        "Avoid N+1 with prefetch_related",
        "Index verification with EXPLAIN ANALYZE",
        "Optional caching: Redis for popular lists",
    ])
    add_footer(s, PROJECT_NAME, 11)

    # 12) Documentation & DX
    s = add_bullets_slide(prs, "Documentation & Developer Experience", [
        "Swagger/Redoc (OpenAPI)",
        "README: setup, run, env, auth flow, endpoints",
        "Postman collection (optional)",
        "Pre-commit: Black, Flake8, isort",
    ])
    add_footer(s, PROJECT_NAME, 12)

    # 13) Testing & Quality
    s = add_bullets_slide(prs, "Testing & Code Quality", [
        "Unit: serializers, utils",
        "Integration: endpoints (auth, products, categories)",
        "Coverage target: 80%+",
        "CI: tests, lint, migrations check",
    ])
    add_footer(s, PROJECT_NAME, 13)

    # 14) Deployment
    s = add_bullets_slide(prs, "Deployment", [
        "Hosting: Platform + Gunicorn",
        "Environment: secrets via env vars",
        "Health: /health (200 OK)",
        "Run migrations on release",
    ])
    add_footer(s, PROJECT_NAME, 14)

    # 15) Live Demo Flow
    s = add_bullets_slide(prs, "Live Demo (≤5 minutes)", [
        "Login: get JWT tokens",
        "Create Category; Create Product (with category)",
        "List Products: filter, sort, paginate",
        "Update/Delete Product (permission gate)",
    ])
    add_footer(s, PROJECT_NAME, 15)

    # 16) Rubric Mapping
    s = add_bullets_slide(prs, "Rubric Mapping", [
        "Functionality: CRUD + auth + discovery",
        "Code Quality: readable, docs, linters",
        "Design & API: normalized ERD, RESTful endpoints, DRF ORM",
        "Deployment & Best Practices: live, secure, documented",
    ])
    add_footer(s, PROJECT_NAME, 16)

    # 17) Challenges & Learnings
    s = add_bullets_slide(prs, "Challenges & Learnings", [
        "Pagination + ordering performance",
        "Auth edge cases and JWT rotation",
        "EXPLAIN plans and index choice",
        "OpenAPI workflows",
    ])
    add_footer(s, PROJECT_NAME, 17)

    # 18) Roadmap
    s = add_bullets_slide(prs, "Roadmap", [
        "Features: search, inventory, orders/cart, reviews",
        "Performance: Redis cache, selective indexes",
        "Security: 2FA, admin audit logs, rate limits",
        "Platform: staging, blue/green, observability",
    ])
    add_footer(s, PROJECT_NAME, 18)

    # 19) Links & Access
    s = add_bullets_slide(prs, "Links & Access", [
        f"Repo: {REPO_URL}",
        "Live API: [URL]",
        "Swagger/Redoc: [URL]",
        "ERD (Google Doc): [URL – anyone with link can view]",
        "Demo video (≤5 min): [URL]",
        "Slide deck: [Google Slides link]",
        "Health: [URL]/health",
    ])
    add_footer(s, PROJECT_NAME, 19)

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("/workspace/alx_project_nexus_ecommerce.pptx")